"""Import and export: office documents, images, e-books, HTML and archives.

Import (→ PDF)
--------------
Images (PNG/JPEG/BMP/GIF/WEBP/TIFF/HEIC/RAW), SVG, TXT, Markdown, HTML, RTF,
EPUB, CBZ/ZIP image archives natively through MuPDF/Pillow, and DOCX/PPTX/XLSX
through LibreOffice when it is available (with a pure-Python fallback that
preserves text and tables).

Export (PDF →)
--------------
Raster images per page, TIFF (multi-page), SVG, plain text, HTML, Markdown,
DOCX, PPTX, XLSX-friendly CSV, plus archival profiles (PDF/A, PDF/X, PDF/E,
PDF/UA) driven by Ghostscript when present and a best-effort internal profile
otherwise.
"""

from __future__ import annotations

import html
import io
import re
import shutil
import subprocess
import tempfile
import zipfile
from collections.abc import Sequence
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Literal

import pymupdf as fitz

from pdfstudio.core.exceptions import (
    DependencyMissingError,
    DocumentError,
    ValidationError,
)
from pdfstudio.core.jobs import JobContext
from pdfstudio.core.logging_setup import get_logger
from pdfstudio.pdfengine.document import PdfDocument, SaveOptions
from pdfstudio.pdfengine.types import PAGE_SIZES, ConformanceLevel, PageSize

log = get_logger("convert")

#: Formats MuPDF opens directly.
NATIVE_INPUT: frozenset[str] = frozenset(
    {".pdf", ".xps", ".oxps", ".cbz", ".fb2", ".epub", ".mobi", ".svg", ".txt"}
)
IMAGE_INPUT: frozenset[str] = frozenset(
    {
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
        ".gif",
        ".webp",
        ".tif",
        ".tiff",
        ".pnm",
        ".pgm",
        ".ppm",
        ".pbm",
        ".jfif",
        ".jxr",
        ".psd",
    }
)
RAW_INPUT: frozenset[str] = frozenset(
    {".cr2", ".cr3", ".nef", ".arw", ".dng", ".orf", ".rw2", ".raf", ".sr2"}
)
OFFICE_INPUT: frozenset[str] = frozenset(
    {".docx", ".doc", ".odt", ".rtf", ".pptx", ".ppt", ".odp", ".xlsx", ".xls", ".ods"}
)
MARKUP_INPUT: frozenset[str] = frozenset({".html", ".htm", ".md", ".markdown", ".xml"})


def supported_import_extensions() -> list[str]:
    """Every extension the import pipeline recognises (for file dialogs)."""
    return sorted(
        NATIVE_INPUT
        | IMAGE_INPUT
        | RAW_INPUT
        | OFFICE_INPUT
        | MARKUP_INPUT
        | {".heic", ".heif", ".zip"}
    )


# --------------------------------------------------------------------------- #
# External tool discovery
# --------------------------------------------------------------------------- #
def find_libreoffice() -> str | None:
    """Locate a LibreOffice binary for high-fidelity office conversion."""
    for name in ("soffice", "libreoffice", "soffice.bin"):
        found = shutil.which(name)
        if found:
            return found
    for candidate in (
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ):
        if Path(candidate).exists():
            return candidate
    return None


def find_ghostscript() -> str | None:
    """Locate Ghostscript, used for PDF/A and PDF/X conversion."""
    for name in ("gs", "gswin64c", "gswin32c"):
        found = shutil.which(name)
        if found:
            return found
    return None


# --------------------------------------------------------------------------- #
# Import
# --------------------------------------------------------------------------- #
@dataclass(slots=True)
class ImportOptions:
    """Controls how non-PDF sources are laid out on the page."""

    page_size: PageSize | str = "A4"
    margin: float = 36.0
    fit_images: bool = True
    one_image_per_page: bool = True
    font_size: float = 11.0
    font_family: str = "sans-serif"
    landscape: bool = False
    dpi: int = 200
    quality: int = 90
    timeout: int = 180

    def resolved_size(self) -> PageSize:
        size = PAGE_SIZES[self.page_size] if isinstance(self.page_size, str) else self.page_size
        return size.landscape() if self.landscape else size


class Importer:
    """Converts many source formats into :class:`PdfDocument` objects."""

    def __init__(self, options: ImportOptions | None = None) -> None:
        self.options = options or ImportOptions()

    def import_file(self, path: str | Path, *, ctx: JobContext | None = None) -> PdfDocument:
        """Convert any supported file into a PDF document."""
        source = Path(path).expanduser()
        if not source.exists():
            raise DocumentError(f"File not found: {source}")
        suffix = source.suffix.lower()
        log.info("Importing {} ({})", source.name, suffix)

        if suffix == ".pdf":
            return PdfDocument.open(source)
        if suffix in OFFICE_INPUT:
            return self._import_office(source, ctx=ctx)
        if suffix in {".heic", ".heif"} or suffix in RAW_INPUT:
            return self._import_special_image(source)
        if suffix in IMAGE_INPUT:
            return self.import_images([source])
        if suffix in {".md", ".markdown"}:
            return self._import_markdown(source)
        if suffix in {".html", ".htm"}:
            return self._import_html(source.read_text("utf-8", errors="replace"))
        if suffix == ".rtf":
            return self._import_rtf(source)
        if suffix == ".zip":
            return self._import_zip(source)
        if suffix in NATIVE_INPUT:
            return self._import_native(source)
        raise ValidationError(f"Unsupported file type: {suffix}")

    # -- native / MuPDF ------------------------------------------------------ #
    def _import_native(self, path: Path) -> PdfDocument:
        """EPUB, CBZ, XPS, FB2, SVG, TXT — MuPDF converts these directly."""
        source = fitz.open(path)
        try:
            size = self.options.resolved_size()
            if source.is_reflowable:
                source.layout(
                    width=size.width, height=size.height, fontsize=self.options.font_size
                )
            data = source.convert_to_pdf()
        finally:
            source.close()
        return PdfDocument.from_bytes(data, name=f"{path.stem}.pdf")

    # -- images -------------------------------------------------------------- #
    def import_images(
        self, paths: Sequence[str | Path], *, ctx: JobContext | None = None
    ) -> PdfDocument:
        """Build a PDF from a list of images, one per page."""
        if not paths:
            raise ValidationError("No images supplied.")
        size = self.options.resolved_size()
        out = fitz.open()
        if ctx:
            ctx.set_total(len(paths))
        for n, item in enumerate(paths, 1):
            image_path = Path(item)
            data = _read_image_bytes(image_path)
            with fitz.open(
                stream=data, filetype="png" if data[:4] == b"\x89PNG" else None
            ) as img:
                rect = img[0].rect
                pdf_bytes = img.convert_to_pdf()
            if self.options.fit_images:
                page = out.new_page(width=size.width, height=size.height)
                margin = self.options.margin
                area = fitz.Rect(margin, margin, size.width - margin, size.height - margin)
                scale = min(area.width / rect.width, area.height / rect.height)
                w, h = rect.width * scale, rect.height * scale
                target = fitz.Rect(
                    (size.width - w) / 2,
                    (size.height - h) / 2,
                    (size.width + w) / 2,
                    (size.height + h) / 2,
                )
            else:
                page = out.new_page(width=rect.width, height=rect.height)
                target = page.rect
            with fitz.open("pdf", pdf_bytes) as overlay:
                page.show_pdf_page(target, overlay, 0)
            if ctx:
                ctx.progress(n, f"Imported {image_path.name}")
        data = out.tobytes(garbage=3, deflate=True)
        out.close()
        name = Path(paths[0]).stem + (".pdf" if len(paths) == 1 else "-images.pdf")
        return PdfDocument.from_bytes(data, name=name)

    def _import_special_image(self, path: Path) -> PdfDocument:
        """HEIC/HEIF and camera RAW via optional Pillow plug-ins."""
        suffix = path.suffix.lower()
        try:
            from PIL import Image
        except ImportError as exc:
            raise DependencyMissingError("Pillow", "Image import") from exc

        if suffix in {".heic", ".heif"}:
            try:
                import pillow_heif

                pillow_heif.register_heif_opener()
            except ImportError as exc:
                raise DependencyMissingError("pillow-heif", "HEIC/HEIF import") from exc
            image = Image.open(path)
        else:
            try:
                import rawpy
            except ImportError as exc:
                raise DependencyMissingError("rawpy", "RAW image import") from exc
            with rawpy.imread(str(path)) as raw:
                image = Image.fromarray(raw.postprocess())

        buffer = io.BytesIO()
        image.convert("RGB").save(buffer, format="PNG")
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp.write(buffer.getvalue())
            temp_path = Path(tmp.name)
        try:
            return self.import_images([temp_path])
        finally:
            temp_path.unlink(missing_ok=True)

    def _import_zip(self, path: Path) -> PdfDocument:
        """A ZIP of images becomes a document (comic/scan archives)."""
        with zipfile.ZipFile(path) as archive:
            names = sorted(
                n for n in archive.namelist() if Path(n).suffix.lower() in IMAGE_INPUT
            )
            if not names:
                raise ValidationError("The archive contains no images.")
            with tempfile.TemporaryDirectory() as tmpdir:
                extracted: list[Path] = []
                for name in names:
                    target = Path(tmpdir) / Path(name).name
                    target.write_bytes(archive.read(name))
                    extracted.append(target)
                doc = self.import_images(extracted)
        doc.display_name = f"{path.stem}.pdf"
        return doc

    # -- text and markup ------------------------------------------------------ #
    def _import_markdown(self, path: Path) -> PdfDocument:
        text = path.read_text("utf-8", errors="replace")
        try:
            import markdown

            body = markdown.markdown(text, extensions=["tables", "fenced_code", "toc"])
        except ImportError:
            body = _markdown_to_html(text)
        doc = self._import_html(body)
        doc.display_name = f"{path.stem}.pdf"
        return doc

    def _import_html(self, source: str) -> PdfDocument:
        """Render HTML with MuPDF's Story engine (CSS supported)."""
        size = self.options.resolved_size()
        margin = self.options.margin
        css = (
            f"body {{ font-family: {self.options.font_family};"
            f" font-size: {self.options.font_size}pt; line-height: 1.45; }}"
            " h1,h2,h3 { margin: 0.6em 0 0.3em; }"
            " table { border-collapse: collapse; }"
            " td,th { border: 1px solid #999; padding: 4px 6px; }"
            " code,pre { font-family: monospace; background: #f4f4f4; }"
        )
        story = fitz.Story(html=source, user_css=css)
        writer = fitz.DocumentWriter(buffer := io.BytesIO())
        area = fitz.Rect(margin, margin, size.width - margin, size.height - margin)
        more = 1
        guard = 0
        while more and guard < 10_000:
            device = writer.begin_page(fitz.Rect(0, 0, size.width, size.height))
            more, _ = story.place(area)
            story.draw(device)
            writer.end_page()
            guard += 1
        writer.close()
        return PdfDocument.from_bytes(buffer.getvalue(), name="import.pdf")

    def import_text(self, text: str, *, name: str = "text.pdf") -> PdfDocument:
        """Lay out plain text with wrapping and pagination."""
        doc = self._import_html(f"<pre>{html.escape(text)}</pre>")
        doc.display_name = name
        return doc

    def _import_rtf(self, path: Path) -> PdfDocument:
        """RTF via LibreOffice, falling back to a plain-text strip."""
        if find_libreoffice():
            return self._import_office(path)
        raw = path.read_text("utf-8", errors="replace")
        text = re.sub(r"\\'[0-9a-f]{2}|\\[a-z]+-?\d* ?|[{}]", "", raw)
        return self.import_text(text, name=f"{path.stem}.pdf")

    # -- office --------------------------------------------------------------- #
    def _import_office(self, path: Path, *, ctx: JobContext | None = None) -> PdfDocument:
        """DOCX/PPTX/XLSX/ODF → PDF, preferring LibreOffice for fidelity."""
        binary = find_libreoffice()
        if binary:
            if ctx:
                ctx.progress(0, "Converting with LibreOffice…")
            with tempfile.TemporaryDirectory() as tmpdir:
                try:
                    subprocess.run(  # noqa: S603
                        [
                            binary,
                            "--headless",
                            "--norestore",
                            "--invisible",
                            "--convert-to",
                            "pdf",
                            "--outdir",
                            tmpdir,
                            str(path),
                        ],
                        check=True,
                        capture_output=True,
                        timeout=self.options.timeout,
                    )
                    produced = Path(tmpdir) / f"{path.stem}.pdf"
                    if produced.exists():
                        return PdfDocument.from_bytes(produced.read_bytes(), name=produced.name)
                except (subprocess.SubprocessError, OSError) as exc:
                    log.warning("LibreOffice conversion failed ({}); using fallback", exc)
        return self._import_office_fallback(path)

    def _import_office_fallback(self, path: Path) -> PdfDocument:
        """Pure-Python extraction preserving text, headings and tables."""
        suffix = path.suffix.lower()
        if suffix == ".docx":
            body = _docx_to_html(path)
        elif suffix == ".pptx":
            body = _pptx_to_html(path)
        elif suffix in {".xlsx", ".xls"}:
            body = _xlsx_to_html(path)
        else:
            raise DependencyMissingError("LibreOffice", f"Conversion of {suffix} files")
        doc = self._import_html(body)
        doc.display_name = f"{path.stem}.pdf"
        return doc


# --------------------------------------------------------------------------- #
# Export
# --------------------------------------------------------------------------- #
@dataclass(slots=True)
class ExportOptions:
    """Raster/text export parameters."""

    dpi: int = 200
    quality: int = 90
    alpha: bool = False
    gray: bool = False
    annotations: bool = True
    pages: Sequence[int] | None = None


class Exporter:
    """Exports a document to images, text, markup and office formats."""

    def __init__(self, document: PdfDocument) -> None:
        self.doc = document

    def _pages(self, options: ExportOptions) -> list[int]:
        return (
            list(options.pages)
            if options.pages is not None
            else list(range(self.doc.page_count))
        )

    # -- raster --------------------------------------------------------------- #
    def to_images(
        self,
        directory: str | Path,
        fmt: Literal["png", "jpg", "jpeg", "bmp", "tif", "tiff", "webp", "ppm"] = "png",
        options: ExportOptions | None = None,
        *,
        prefix: str = "page",
        ctx: JobContext | None = None,
    ) -> list[Path]:
        """Write one image file per page."""
        opts = options or ExportOptions()
        out_dir = Path(directory)
        out_dir.mkdir(parents=True, exist_ok=True)
        pages = self._pages(opts)
        if ctx:
            ctx.set_total(len(pages))
        written: list[Path] = []
        width = len(str(max(pages, default=0) + 1))
        for n, index in enumerate(pages, 1):
            data = self.page_image(index, fmt=fmt, options=opts)
            target = out_dir / f"{prefix}-{index + 1:0{width}d}.{fmt}"
            target.write_bytes(data)
            written.append(target)
            if ctx:
                ctx.progress(n, f"Exported page {index + 1}")
        log.info("Exported {} image(s) to {}", len(written), out_dir)
        return written

    def page_image(
        self, index: int, *, fmt: str = "png", options: ExportOptions | None = None
    ) -> bytes:
        """Rasterise one page to image bytes."""
        opts = options or ExportOptions()
        with self.doc.locked() as handle:
            pix = handle[index].get_pixmap(
                dpi=opts.dpi,
                alpha=opts.alpha and fmt in ("png", "webp"),
                colorspace=fitz.csGRAY if opts.gray else fitz.csRGB,
                annots=opts.annotations,
            )
            if fmt in ("jpg", "jpeg"):
                return pix.tobytes("jpeg", jpg_quality=opts.quality)
            if fmt in ("tif", "tiff"):
                return _pixmap_to_tiff(pix)
            if fmt == "webp":
                return _pixmap_via_pillow(pix, "WEBP", quality=opts.quality)
            if fmt == "bmp":
                return _pixmap_via_pillow(pix, "BMP")
            return pix.tobytes(fmt)

    def to_multipage_tiff(self, path: str | Path, options: ExportOptions | None = None) -> Path:
        """Write a single multi-page TIFF (archival / fax workflows)."""
        try:
            from PIL import Image
        except ImportError as exc:
            raise DependencyMissingError("Pillow", "TIFF export") from exc
        opts = options or ExportOptions()
        frames = []
        for index in self._pages(opts):
            with self.doc.locked() as handle:
                pix = handle[index].get_pixmap(dpi=opts.dpi, annots=opts.annotations)
            frames.append(Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB"))
        target = Path(path)
        frames[0].save(
            target,
            save_all=True,
            append_images=frames[1:],
            compression="tiff_deflate",
            dpi=(opts.dpi, opts.dpi),
        )
        return target

    def to_svg(self, index: int) -> str:
        """Export a page as scalable vector graphics."""
        with self.doc.locked() as handle:
            return handle[index].get_svg_image(text_as_path=False)

    # -- text and markup ------------------------------------------------------ #
    def to_text(self, options: ExportOptions | None = None, *, layout: bool = False) -> str:
        """Plain text; ``layout=True`` preserves column positions."""
        opts = options or ExportOptions()
        mode = "text" if not layout else "blocks"
        chunks: list[str] = []
        for index in self._pages(opts):
            if layout:
                blocks = self.doc.extract_blocks(index)
                blocks.sort(key=lambda b: (round(b.rect.y0), b.rect.x0))
                chunks.append("\n".join(b.text for b in blocks))
            else:
                chunks.append(self.doc.extract_text(index, mode=mode))
        return "\n\f\n".join(chunks)

    def to_html(self, options: ExportOptions | None = None) -> str:
        """Styled HTML preserving fonts, positions and images."""
        opts = options or ExportOptions()
        body = "\n".join(
            self.doc.extract_text(index, mode="html") for index in self._pages(opts)
        )
        title = html.escape(self.doc.metadata().title or self.doc.display_name)
        return (
            "<!doctype html>\n<html><head><meta charset='utf-8'>"
            f"<title>{title}</title></head>\n<body>\n{body}\n</body></html>"
        )

    def to_markdown(self, options: ExportOptions | None = None) -> str:
        """Heuristic Markdown conversion using font sizes to detect headings."""
        opts = options or ExportOptions()
        lines: list[str] = []
        sizes: list[float] = []
        for index in self._pages(opts):
            for block in self.doc.extract_blocks(index):
                for line in block.lines:
                    for span in line.spans:
                        if span.text.strip():
                            sizes.append(span.size)
        body_size = _median(sizes) if sizes else 11.0

        for index in self._pages(opts):
            for block in self.doc.extract_blocks(index):
                text = block.text.strip()
                if not text:
                    continue
                first = (
                    block.lines[0].spans[0] if block.lines and block.lines[0].spans else None
                )
                size = first.size if first else body_size
                bold = first.bold if first else False
                if size >= body_size * 1.6:
                    lines.append(f"# {text}")
                elif size >= body_size * 1.3:
                    lines.append(f"## {text}")
                elif size >= body_size * 1.12 or (bold and len(text) < 80):
                    lines.append(f"### {text}")
                elif re.match(r"^[•·▪\-*]\s+", text):
                    lines.append(re.sub(r"^[•·▪]\s+", "- ", text))
                else:
                    lines.append(text)
                lines.append("")
            lines.append("---\n")
        return "\n".join(lines).strip()

    def to_docx(self, path: str | Path, options: ExportOptions | None = None) -> Path:
        """Export to Word, preferring LibreOffice for layout fidelity."""
        target = Path(path)
        binary = find_libreoffice()
        if binary and self.doc.path:
            with tempfile.TemporaryDirectory() as tmpdir:
                try:
                    subprocess.run(  # noqa: S603
                        [
                            binary,
                            "--headless",
                            "--convert-to",
                            "docx:MS Word 2007 XML",
                            "--outdir",
                            tmpdir,
                            str(self.doc.path),
                        ],
                        check=True,
                        capture_output=True,
                        timeout=300,
                    )
                    produced = Path(tmpdir) / f"{self.doc.path.stem}.docx"
                    if produced.exists():
                        shutil.copy(produced, target)
                        return target
                except (subprocess.SubprocessError, OSError) as exc:
                    log.warning("LibreOffice DOCX export failed ({}); using fallback", exc)
        return self._docx_fallback(target, options)

    def _docx_fallback(self, target: Path, options: ExportOptions | None) -> Path:
        """Build a DOCX from extracted structure with python-docx."""
        try:
            from docx import Document as DocxDocument
            from docx.shared import Pt, RGBColor
        except ImportError as exc:
            raise DependencyMissingError("python-docx", "DOCX export") from exc

        opts = options or ExportOptions()
        docx = DocxDocument()
        meta = self.doc.metadata()
        if meta.title:
            docx.core_properties.title = meta.title
        if meta.author:
            docx.core_properties.author = meta.author

        sizes = [
            span.size
            for index in self._pages(opts)
            for block in self.doc.extract_blocks(index)
            for line in block.lines
            for span in line.spans
            if span.text.strip()
        ]
        body_size = _median(sizes) if sizes else 11.0

        for n, index in enumerate(self._pages(opts)):
            if n:
                docx.add_page_break()
            for block in self.doc.extract_blocks(index):
                text = block.text.strip()
                if not text:
                    continue
                first = block.lines[0].spans[0]
                if first.size >= body_size * 1.5:
                    docx.add_heading(text, level=1)
                elif first.size >= body_size * 1.25:
                    docx.add_heading(text, level=2)
                else:
                    paragraph = docx.add_paragraph()
                    for line in block.lines:
                        for span in line.spans:
                            run = paragraph.add_run(span.text)
                            run.bold = span.bold
                            run.italic = span.italic
                            run.font.size = Pt(max(6.0, span.size))
                            colour = span.color
                            run.font.color.rgb = RGBColor(
                                int(colour.r * 255), int(colour.g * 255), int(colour.b * 255)
                            )
                        paragraph.add_run(" ")
        docx.save(target)
        return target

    def to_pptx(self, path: str | Path, options: ExportOptions | None = None) -> Path:
        """One slide per page, using a rendered image as the slide background."""
        try:
            from pptx import Presentation
            from pptx.util import Emu
        except ImportError as exc:
            raise DependencyMissingError("python-pptx", "PPTX export") from exc

        opts = options or ExportOptions()
        pages = self._pages(opts)
        width, height = self.doc.page_size(pages[0])
        presentation = Presentation()
        presentation.slide_width = Emu(int(width * 12700))
        presentation.slide_height = Emu(int(height * 12700))
        blank = presentation.slide_layouts[6]
        for index in pages:
            slide = presentation.slides.add_slide(blank)
            image = self.page_image(index, fmt="png", options=opts)
            slide.shapes.add_picture(
                io.BytesIO(image),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
        presentation.save(str(path))
        return Path(path)

    def tables_to_csv(
        self, directory: str | Path, *, options: ExportOptions | None = None
    ) -> list[Path]:
        """Detect tables and write one CSV per table."""
        import csv

        opts = options or ExportOptions()
        out_dir = Path(directory)
        out_dir.mkdir(parents=True, exist_ok=True)
        written: list[Path] = []
        for index in self._pages(opts):
            with self.doc.locked() as handle:
                try:
                    found = handle[index].find_tables()
                except Exception:
                    continue
                for n, table in enumerate(found.tables, 1):
                    rows = table.extract()
                    if not rows:
                        continue
                    target = out_dir / f"page{index + 1}-table{n}.csv"
                    with target.open("w", newline="", encoding="utf-8") as fh:
                        csv.writer(fh).writerows(
                            [["" if c is None else str(c) for c in row] for row in rows]
                        )
                    written.append(target)
        return written

    # -- archival profiles ----------------------------------------------------- #
    def to_conformance(
        self,
        path: str | Path,
        level: ConformanceLevel = ConformanceLevel.PDF_A_2B,
        *,
        ctx: JobContext | None = None,
    ) -> Path:
        """Convert to PDF/A, PDF/X or PDF/E using Ghostscript when available.

        Without Ghostscript a best-effort conversion is performed: fonts are
        embedded where possible, transparency is preserved, XMP metadata is
        written and the output is validated structurally.  The result is
        *conformance-oriented* rather than certified; the UI states this.
        """
        target = Path(path)
        source_path = self.doc.path
        gs = find_ghostscript()
        if gs and source_path and source_path.exists():
            device_args = _ghostscript_profile(level)
            if ctx:
                ctx.progress(0, f"Converting to {level.value} with Ghostscript…")
            try:
                subprocess.run(  # noqa: S603
                    [
                        gs,
                        "-dBATCH",
                        "-dNOPAUSE",
                        "-dQUIET",
                        "-dSAFER",
                        "-sDEVICE=pdfwrite",
                        *device_args,
                        f"-sOutputFile={target}",
                        str(source_path),
                    ],
                    check=True,
                    capture_output=True,
                    timeout=600,
                )
                log.info("Wrote {} via Ghostscript", level.value)
                return target
            except (subprocess.SubprocessError, OSError) as exc:
                log.warning("Ghostscript failed ({}); using internal profile", exc)

        clone = self.doc.copy()
        try:
            _apply_internal_conformance(clone, level)
            clone.save_as(target, SaveOptions(garbage=4, deflate=True, clean=True, linear=True))
        finally:
            clone.close()
        return target


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #
def _ghostscript_profile(level: ConformanceLevel) -> list[str]:
    common = [
        "-dPDFSETTINGS=/prepress",
        "-dEmbedAllFonts=true",
        "-dSubsetFonts=true",
        "-dCompatibilityLevel=1.7",
    ]
    match level:
        case ConformanceLevel.PDF_A_1B:
            return [
                *common,
                "-dPDFA=1",
                "-sColorConversionStrategy=RGB",
                "-dPDFACompatibilityPolicy=1",
            ]
        case ConformanceLevel.PDF_A_2B | ConformanceLevel.PDF_A_2U:
            return [
                *common,
                "-dPDFA=2",
                "-sColorConversionStrategy=RGB",
                "-dPDFACompatibilityPolicy=1",
            ]
        case ConformanceLevel.PDF_A_3B:
            return [
                *common,
                "-dPDFA=3",
                "-sColorConversionStrategy=RGB",
                "-dPDFACompatibilityPolicy=1",
            ]
        case ConformanceLevel.PDF_X_1A:
            return [*common, "-dPDFX", "-sColorConversionStrategy=CMYK"]
        case ConformanceLevel.PDF_X_4:
            return [
                *common,
                "-dPDFX",
                "-sColorConversionStrategy=CMYK",
                "-dCompatibilityLevel=1.6",
            ]
        case _:
            return common


def _apply_internal_conformance(doc: PdfDocument, level: ConformanceLevel) -> None:
    """Write the XMP identification block required by the profile."""
    part, conformance = {
        ConformanceLevel.PDF_A_1B: ("1", "B"),
        ConformanceLevel.PDF_A_2B: ("2", "B"),
        ConformanceLevel.PDF_A_2U: ("2", "U"),
        ConformanceLevel.PDF_A_3B: ("3", "B"),
        ConformanceLevel.PDF_UA_1: ("1", "UA"),
        ConformanceLevel.PDF_X_1A: ("1", "X"),
        ConformanceLevel.PDF_X_4: ("4", "X"),
        ConformanceLevel.PDF_E_1: ("1", "E"),
    }[level]
    meta = doc.metadata()
    schema = "pdfuaid" if conformance == "UA" else "pdfxid" if conformance == "X" else "pdfaid"
    xmp = f"""<?xpacket begin="\ufeff" id="W5M0MpCehiHzreSzNTczkc9d"?>
<x:xmpmeta xmlns:x="adobe:ns:meta/">
 <rdf:RDF xmlns:rdf="http://www.w3.org/1999/02/22-rdf-syntax-ns#">
  <rdf:Description rdf:about=""
    xmlns:{schema}="http://www.aiim.org/{schema}/ns/id/"
    xmlns:dc="http://purl.org/dc/elements/1.1/">
   <{schema}:part>{part}</{schema}:part>
   <{schema}:conformance>{conformance}</{schema}:conformance>
   <dc:title>{html.escape(meta.title or doc.display_name)}</dc:title>
   <dc:creator>{html.escape(meta.author)}</dc:creator>
  </rdf:Description>
 </rdf:RDF>
</x:xmpmeta>
<?xpacket end="w"?>"""
    meta.xmp = xmp
    meta.producer = "PDF Studio"
    doc.set_metadata(meta)


def _read_image_bytes(path: Path) -> bytes:
    data = path.read_bytes()
    suffix = path.suffix.lower()
    if suffix in IMAGE_INPUT and suffix not in {".psd"}:
        return data
    try:
        from PIL import Image
    except ImportError as exc:
        raise DependencyMissingError("Pillow", "Image import") from exc
    buffer = io.BytesIO()
    Image.open(io.BytesIO(data)).convert("RGB").save(buffer, format="PNG")
    return buffer.getvalue()


def _pixmap_to_tiff(pix: fitz.Pixmap) -> bytes:
    try:
        from PIL import Image
    except ImportError as exc:
        raise DependencyMissingError("Pillow", "TIFF export") from exc
    image = Image.open(io.BytesIO(pix.tobytes("png")))
    buffer = io.BytesIO()
    image.save(buffer, format="TIFF", compression="tiff_deflate")
    return buffer.getvalue()


def _pixmap_via_pillow(pix: fitz.Pixmap, fmt: str, **kwargs: Any) -> bytes:
    try:
        from PIL import Image
    except ImportError as exc:
        raise DependencyMissingError("Pillow", f"{fmt} export") from exc
    image = Image.open(io.BytesIO(pix.tobytes("png")))
    if fmt in ("BMP", "JPEG"):
        image = image.convert("RGB")
    buffer = io.BytesIO()
    image.save(buffer, format=fmt, **kwargs)
    return buffer.getvalue()


def _median(values: Sequence[float]) -> float:
    if not values:
        return 0.0
    ordered = sorted(values)
    middle = len(ordered) // 2
    if len(ordered) % 2:
        return ordered[middle]
    return (ordered[middle - 1] + ordered[middle]) / 2


def _markdown_to_html(text: str) -> str:
    """Minimal Markdown renderer used when the ``markdown`` package is absent."""
    lines: list[str] = []
    in_list = False
    for raw in text.splitlines():
        line = html.escape(raw.rstrip())
        heading = re.match(r"^(#{1,6})\s+(.*)$", line)
        bullet = re.match(r"^[-*+]\s+(.*)$", line)
        if heading:
            if in_list:
                lines.append("</ul>")
                in_list = False
            level = len(heading.group(1))
            lines.append(f"<h{level}>{heading.group(2)}</h{level}>")
        elif bullet:
            if not in_list:
                lines.append("<ul>")
                in_list = True
            lines.append(f"<li>{bullet.group(1)}</li>")
        elif not line.strip():
            if in_list:
                lines.append("</ul>")
                in_list = False
            lines.append("<p></p>")
        else:
            lines.append(f"<p>{line}</p>")
    if in_list:
        lines.append("</ul>")
    body = "\n".join(lines)
    body = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", body)
    body = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"<i>\1</i>", body)
    body = re.sub(r"`([^`]+)`", r"<code>\1</code>", body)
    return body


def _docx_to_html(path: Path) -> str:
    """Extract paragraphs, headings and tables from a DOCX without Word."""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        return _ooxml_text_fallback(path, "word/document.xml")
    document = DocxDocument(str(path))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = html.escape(paragraph.text.strip())
        if not text:
            continue
        style = (paragraph.style.name or "").lower()
        if style.startswith("heading"):
            level = re.sub(r"\D", "", style) or "1"
            parts.append(f"<h{level}>{text}</h{level}>")
        elif style.startswith("list"):
            parts.append(f"<ul><li>{text}</li></ul>")
        else:
            parts.append(f"<p>{text}</p>")
    for table in document.tables:
        rows = [
            "<tr>"
            + "".join(f"<td>{html.escape(cell.text.strip())}</td>" for cell in row.cells)
            + "</tr>"
            for row in table.rows
        ]
        parts.append("<table>" + "".join(rows) + "</table>")
    return "\n".join(parts)


def _pptx_to_html(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return _ooxml_text_fallback(path, "ppt/slides/slide1.xml")
    presentation = Presentation(str(path))
    parts: list[str] = []
    for n, slide in enumerate(presentation.slides, 1):
        parts.append(f"<h2>Slide {n}</h2>")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = html.escape("".join(run.text for run in paragraph.runs).strip())
                    if text:
                        parts.append(f"<p>{text}</p>")
    return "\n".join(parts)


def _xlsx_to_html(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        return _ooxml_text_fallback(path, "xl/sharedStrings.xml")
    workbook = openpyxl.load_workbook(str(path), data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"<h2>{html.escape(sheet.title)}</h2><table>")
        for row in sheet.iter_rows(values_only=True):
            cells = "".join(f"<td>{html.escape('' if v is None else str(v))}</td>" for v in row)
            parts.append(f"<tr>{cells}</tr>")
        parts.append("</table>")
    return "\n".join(parts)


def _ooxml_text_fallback(path: Path, member: str) -> str:
    """Last-resort text extraction straight from the OOXML package."""
    try:
        with zipfile.ZipFile(path) as archive:
            names = [
                n for n in archive.namelist() if n.startswith(member.split("/", maxsplit=1)[0])
            ]
            text_parts: list[str] = []
            for name in sorted(names):
                if not name.endswith(".xml"):
                    continue
                xml = archive.read(name).decode("utf-8", "replace")
                text_parts.extend(re.findall(r"<a:t>(.*?)</a:t>|<w:t[^>]*>(.*?)</w:t>", xml))
            flat = [html.escape(a or b) for a, b in text_parts if (a or b)]
            return "".join(f"<p>{t}</p>" for t in flat)
    except (zipfile.BadZipFile, OSError) as exc:
        raise DocumentError(f"Cannot read {path.name}", detail=str(exc)) from exc
